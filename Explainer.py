import asyncio
import json
import os
from time import sleep
from pptx import Presentation
from openai import AsyncOpenAI

# Ensure you have set the OpenAI API key in your environment variables
client = AsyncOpenAI(api_key=os.environ.get("OPENAI_API_KEY"))

UPLOADS_FOLDER = r'C:\Users\Naama\Desktop\studies-material\שנה ד\סמסטר ב\אקסלנטים\f2\uploads'
OUTPUTS_FOLDER = r'C:\Users\Naama\Desktop\studies-material\שנה ד\סמסטר ב\אקסלנטים\f2\output'
PROCESSED_FILES_LOG = 'processed_files.log'

# Ensure directories exist
os.makedirs(UPLOADS_FOLDER, exist_ok=True)
os.makedirs(OUTPUTS_FOLDER, exist_ok=True)


async def send_slide_to_openai(slide_text, slide_number, my_prompt):
    try:
        response = await client.chat.completions.create(
            messages=[
                {
                    "role": "user",
                    "content": my_prompt + slide_text,
                }
            ],
            model="gpt-3.5-turbo",
        )
        return {"slide": slide_number, "response": response.choices[0].message.content}
    except Exception as e:
        return {"slide": slide_number, "error": str(e)}


def process_presentation(presentation_path):
    presentation = Presentation(presentation_path)
    tasks = []
    my_prompt = ""
    for i, slide in enumerate(presentation.slides):
        slide_text = ("I am a student, I need your help to understand this slide from a presentation, in your answer, "
                      "simply explain it to me, don't mention its a slide, and act as if I haven't read the presentation yet "
                      "(ie in your answer if the presentation mentions a polygraph for example, in your answer don't say 'this tool' without giving its name): ")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + " "
        # Create a task to send slide text asynchronously
        tasks.append(send_slide_to_openai(slide_text.strip(), i + 1, my_prompt))

    return tasks


async def process_file(file_path, output_path):
    tasks = process_presentation(file_path)
    responses = await asyncio.gather(*tasks)
    with open(output_path, "w") as outfile:
        outfile.write(json.dumps(responses, indent=4))


def load_processed_files():
    if os.path.exists(PROCESSED_FILES_LOG):
        with open(PROCESSED_FILES_LOG, 'r') as f:
            return set(f.read().splitlines())
    return set()


def save_processed_file(filename):
    with open(PROCESSED_FILES_LOG, 'a') as f:
        f.write(filename + '\n')


async def main_loop():
    processed_files = load_processed_files()
    while True:
        for filename in os.listdir(UPLOADS_FOLDER):
            if filename not in processed_files:
                file_path = os.path.join(UPLOADS_FOLDER, filename)
                output_path = os.path.join(OUTPUTS_FOLDER, filename+'.json')

                print(f"Processing file: {filename}")

                await process_file(file_path, output_path)

                print(f"Finished processing file: {filename}")

                processed_files.add(filename)
                save_processed_file(filename)

        sleep(10)  # Sleep for 10 seconds before next iteration


if __name__ == "__main__":
    try:
        asyncio.run(main_loop())
    except KeyboardInterrupt:
        print("Explainer system stopped.")
